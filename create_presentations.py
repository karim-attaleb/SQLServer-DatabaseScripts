#!/usr/bin/env python3
"""
Generate PowerPoint presentations for SQL Server Database Scripts documentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """Create a title slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    return slide

def create_content_slide(prs, title):
    """Create a blank content slide with title"""
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    return slide

def add_bullet_points(slide, left, top, width, height, points):
    """Add bullet points to a slide"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
        p.space_before = Pt(6)
        
    return textbox

def add_process_flow_diagram(slide):
    """Add a process flow diagram showing the database creation workflow"""
    left_margin = Inches(0.8)
    top_start = Inches(1.8)
    box_width = Inches(2.5)
    box_height = Inches(0.8)
    arrow_height = Inches(0.4)
    
    steps = [
        ("Validate Connection", RGBColor(52, 152, 219)),
        ("Create Directories", RGBColor(46, 204, 113)),
        ("Calculate File Count", RGBColor(155, 89, 182)),
        ("Validate Disk Space", RGBColor(230, 126, 34)),
        ("Create Database", RGBColor(231, 76, 60)),
        ("Configure Query Store", RGBColor(26, 188, 156))
    ]
    
    current_top = top_start
    
    for i, (step_text, color) in enumerate(steps):
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left_margin, current_top,
            box_width, box_height
        )
        
        box.fill.solid()
        box.fill.fore_color.rgb = color
        
        box.line.color.rgb = RGBColor(255, 255, 255)
        box.line.width = Pt(2)
        
        text_frame = box.text_frame
        text_frame.text = step_text
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.paragraphs[0].font.size = Pt(16)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        text_frame.vertical_anchor = 1
        
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                left_margin + box_width / 2 - Inches(0.15),
                current_top + box_height,
                Inches(0.3), arrow_height
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(149, 165, 166)
            arrow.line.fill.background()
        
        current_top += box_height + arrow_height

def add_architecture_diagram(slide):
    """Add architecture diagram showing components"""
    
    config_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(2),
        Inches(2.5), Inches(1)
    )
    config_box.fill.solid()
    config_box.fill.fore_color.rgb = RGBColor(241, 196, 15)
    config_box.line.color.rgb = RGBColor(243, 156, 18)
    config_box.text_frame.text = "DatabaseConfig.psd1\n(Configuration)"
    config_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    config_box.text_frame.paragraphs[0].font.size = Pt(14)
    config_box.text_frame.paragraphs[0].font.bold = True
    config_box.text_frame.vertical_anchor = 1
    
    utils_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.5), Inches(2),
        Inches(2.5), Inches(1)
    )
    utils_box.fill.solid()
    utils_box.fill.fore_color.rgb = RGBColor(52, 152, 219)
    utils_box.line.color.rgb = RGBColor(41, 128, 185)
    utils_box.text_frame.text = "DatabaseUtils.psm1\n(Utility Functions)"
    utils_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    utils_box.text_frame.paragraphs[0].font.size = Pt(14)
    utils_box.text_frame.paragraphs[0].font.bold = True
    utils_box.text_frame.vertical_anchor = 1
    
    main_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.5), Inches(2),
        Inches(2.5), Inches(1)
    )
    main_box.fill.solid()
    main_box.fill.fore_color.rgb = RGBColor(46, 204, 113)
    main_box.line.color.rgb = RGBColor(39, 174, 96)
    main_box.text_frame.text = "Invoke-DatabaseCreation.ps1\n(Main Script)"
    main_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    main_box.text_frame.paragraphs[0].font.size = Pt(14)
    main_box.text_frame.paragraphs[0].font.bold = True
    main_box.text_frame.vertical_anchor = 1
    
    arrow1 = slide.shapes.add_connector(
        1, Inches(3), Inches(2.5),
        Inches(3.5), Inches(2.5)
    )
    arrow1.line.color.rgb = RGBColor(127, 140, 141)
    arrow1.line.width = Pt(2)
    
    arrow2 = slide.shapes.add_connector(
        1, Inches(6), Inches(2.5),
        Inches(6.5), Inches(2.5)
    )
    arrow2.line.color.rgb = RGBColor(127, 140, 141)
    arrow2.line.width = Pt(2)
    
    sql_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.5), Inches(4.5),
        Inches(2.5), Inches(1.2)
    )
    sql_box.fill.solid()
    sql_box.fill.fore_color.rgb = RGBColor(231, 76, 60)
    sql_box.line.color.rgb = RGBColor(192, 57, 43)
    sql_box.line.width = Pt(3)
    sql_box.text_frame.text = "🗄️ SQL Server\nDatabase"
    sql_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    sql_box.text_frame.paragraphs[0].font.size = Pt(16)
    sql_box.text_frame.paragraphs[0].font.bold = True
    sql_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    sql_box.text_frame.vertical_anchor = 1
    
    arrow3 = slide.shapes.add_connector(
        1, Inches(7.75), Inches(3),
        Inches(4.75), Inches(4.5)
    )
    arrow3.line.color.rgb = RGBColor(127, 140, 141)
    arrow3.line.width = Pt(2)

def add_features_diagram(slide):
    """Add features diagram with icons"""
    features = [
        ("📚", "Comprehensive\nDocumentation", RGBColor(52, 152, 219)),
        ("🔄", "Auto File\nCalculation", RGBColor(46, 204, 113)),
        ("💾", "Disk Space\nValidation", RGBColor(230, 126, 34)),
        ("🧪", "Pester\nTests", RGBColor(155, 89, 182))
    ]
    
    left_start = Inches(1.2)
    top_pos = Inches(2.2)
    box_width = Inches(1.8)
    box_height = Inches(1.5)
    spacing = Inches(0.4)
    
    for i, (icon, feature_text, color) in enumerate(features):
        left_pos = left_start + (i * (box_width + spacing))
        
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left_pos, top_pos,
            box_width, box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = RGBColor(255, 255, 255)
        box.line.width = Pt(2)
        
        text_frame = box.text_frame
        text_frame.clear()
        
        p1 = text_frame.paragraphs[0]
        p1.text = icon
        p1.alignment = PP_ALIGN.CENTER
        p1.font.size = Pt(40)
        
        p2 = text_frame.add_paragraph()
        p2.text = feature_text
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.space_before = Pt(6)
        
        text_frame.vertical_anchor = 1

def create_overview_presentation():
    """Create the main overview presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    create_title_slide(
        prs,
        "SQL Server Database Scripts",
        "Professional Automation for Database Creation\nby karim-attaleb"
    )
    
    slide = create_content_slide(prs, "📋 Project Overview")
    add_bullet_points(slide, Inches(1), Inches(2), Inches(8), Inches(4.5), [
        "🎯 Automated SQL Server database creation with best practices",
        "📊 Intelligent file distribution based on expected database size",
        "🔒 Pre-validation of disk space to prevent failures",
        "📚 Comprehensive documentation and professional structure",
        "🧪 Complete test coverage with Pester testing framework",
        "⚙️ Configurable via PowerShell data files (PSD1)"
    ])
    
    slide = create_content_slide(prs, "🏗️ Architecture Overview")
    add_architecture_diagram(slide)
    
    slide = create_content_slide(prs, "🔄 Database Creation Workflow")
    add_process_flow_diagram(slide)
    
    slide = create_content_slide(prs, "✨ Key Features")
    add_features_diagram(slide)
    
    slide = create_content_slide(prs, "🔢 Simplified File Count Logic")
    add_bullet_points(slide, Inches(1), Inches(2), Inches(8), Inches(2), [
        "📐 Logic: If ExpectedSize > Threshold → Calculate files, else → 1 file",
        "📊 Formula: Ceiling(ExpectedDatabaseSize / FileSizeThreshold)",
        "📈 Example 1: 5GB database ÷ 10GB threshold = 1 file",
        "📈 Example 2: 50GB database ÷ 10GB threshold = 5 files (max 8)",
        "🎯 All files are created in the PRIMARY filegroup"
    ])
    
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(4.2),
        Inches(7), Inches(1.8)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(236, 240, 241)
    box.line.color.rgb = RGBColor(52, 73, 94)
    box.line.width = Pt(2)
    
    tf = box.text_frame
    tf.text = "ExpectedDatabaseSize = \"50GB\"\nFileSizeThreshold = \"10GB\"\n→ Result: 5 data files created automatically"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.name = "Courier New"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    slide = create_content_slide(prs, "💾 Disk Space Validation")
    add_bullet_points(slide, Inches(1), Inches(2), Inches(4), Inches(4), [
        "✅ Pre-flight check using Get-DbaDiskSpace",
        "📊 Calculates: (Files × Size) + LogSize + 10% margin",
        "🛡️ Prevents out-of-space failures during creation",
        "⚠️ Clear error messages if insufficient space",
        "🔍 Validates both data and log drives separately"
    ])
    
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.5), Inches(2.5),
        Inches(3.5), Inches(3)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(230, 126, 34)
    box.line.color.rgb = RGBColor(211, 84, 0)
    box.line.width = Pt(3)
    
    tf = box.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = "💡 Example"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "\n4 files × 200MB\n+ 100MB log\n= 900MB required\n\n+ 10% margin\n= 990MB needed"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    
    slide = create_content_slide(prs, "🧪 Testing Strategy")
    add_bullet_points(slide, Inches(1), Inches(2), Inches(8), Inches(4), [
        "✅ Comprehensive Pester test suite covering all functions",
        "🔍 Unit tests for utility functions (size conversion, file calculation)",
        "🎯 Integration tests for main database creation workflow",
        "⚠️ Error scenario testing (invalid input, insufficient space)",
        "📊 Edge case coverage (same drive for data/log, maximum files)",
        "🔄 Mocked dependencies for isolated testing"
    ])
    
    slide = create_content_slide(prs, "📦 Module Structure")
    
    left = Inches(1.5)
    top = Inches(2)
    width = Inches(7)
    
    items = [
        ("DatabaseConfig.psd1", "Configuration file with all settings", RGBColor(241, 196, 15)),
        ("DatabaseUtils.psm1", "Utility functions module", RGBColor(52, 152, 219)),
        ("DatabaseUtils.psd1", "Module manifest", RGBColor(52, 152, 219)),
        ("Invoke-DatabaseCreation.ps1", "Main execution script", RGBColor(46, 204, 113)),
        ("Tests/", "Pester test files", RGBColor(155, 89, 182)),
        ("README.md", "Complete documentation", RGBColor(149, 165, 166))
    ]
    
    box_height = Inches(0.6)
    spacing = Inches(0.15)
    
    for i, (name, desc, color) in enumerate(items):
        current_top = top + (i * (box_height + spacing))
        
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, current_top,
            width, box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = RGBColor(255, 255, 255)
        box.line.width = Pt(1)
        
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"{name}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        p2 = tf.add_paragraph()
        p2.text = f"  {desc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(255, 255, 255)
        
        tf.vertical_anchor = 1
    
    slide = create_content_slide(prs, "🚀 Quick Start Guide")
    add_bullet_points(slide, Inches(1), Inches(2), Inches(8), Inches(4.5), [
        "1️⃣ Install dbatools: Install-Module -Name dbatools",
        "2️⃣ Configure DatabaseConfig.psd1 with your settings",
        "3️⃣ Set ExpectedDatabaseSize for auto file calculation (optional)",
        "4️⃣ Run: .\\Invoke-DatabaseCreation.ps1 -ConfigPath .\\DatabaseConfig.psd1",
        "5️⃣ Script validates space, creates database, enables Query Store",
        "6️⃣ Check logs for detailed operation information"
    ])
    
    slide = create_content_slide(prs, "📊 Benefits & Impact")
    
    benefits = [
        ("⚡", "Speed", "Automates manual\nprocess", RGBColor(52, 152, 219)),
        ("✅", "Reliability", "Pre-validation\nprevents errors", RGBColor(46, 204, 113)),
        ("📏", "Best Practices", "Follows SQL Server\nguidelines", RGBColor(155, 89, 182)),
        ("🔧", "Maintainable", "Well-documented\nand tested", RGBColor(230, 126, 34))
    ]
    
    left_start = Inches(1)
    top_pos = Inches(2.5)
    box_width = Inches(2)
    box_height = Inches(2.5)
    spacing = Inches(0.2)
    
    for i, (icon, title, desc, color) in enumerate(benefits):
        left_pos = left_start + (i * (box_width + spacing))
        
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left_pos, top_pos,
            box_width, box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = RGBColor(255, 255, 255)
        box.line.width = Pt(3)
        
        tf = box.text_frame
        tf.clear()
        
        p1 = tf.paragraphs[0]
        p1.text = icon
        p1.alignment = PP_ALIGN.CENTER
        p1.font.size = Pt(48)
        
        p2 = tf.add_paragraph()
        p2.text = title
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.space_before = Pt(12)
        
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.alignment = PP_ALIGN.CENTER
        p3.font.size = Pt(12)
        p3.font.color.rgb = RGBColor(255, 255, 255)
        p3.space_before = Pt(6)
        
        tf.vertical_anchor = 1
    
    slide = create_content_slide(prs, "📝 Summary")
    add_bullet_points(slide, Inches(1), Inches(2), Inches(8), Inches(4), [
        "✨ Professional PowerShell automation for SQL Server databases",
        "🎯 Intelligent features: auto file calculation + disk space validation",
        "🧪 Complete test coverage ensures reliability",
        "📚 Comprehensive documentation for easy adoption",
        "🔧 Easy to configure and maintain",
        "🚀 Ready for production use"
    ])
    
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(5.5),
        Inches(6), Inches(1)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0, 51, 102)
    box.line.color.rgb = RGBColor(255, 255, 255)
    box.line.width = Pt(2)
    
    tf = box.text_frame
    tf.text = "🌟 Thank you! Questions?"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf.vertical_anchor = 1
    
    prs.save('/home/ubuntu/sqlserver-databasescripts/SQL_Server_Database_Scripts_Overview.pptx')
    print("✅ Created: SQL_Server_Database_Scripts_Overview.pptx")

def create_technical_presentation():
    """Create detailed technical documentation presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    create_title_slide(
        prs,
        "Technical Documentation",
        "SQL Server Database Scripts - Deep Dive"
    )
    
    slide = create_content_slide(prs, "🔧 Core Functions")
    add_bullet_points(slide, Inches(1), Inches(2), Inches(8), Inches(4.5), [
        "Convert-SizeToInt: Converts size strings (MB/GB/TB) to integers",
        "Calculate-OptimalDataFiles: Determines file count from size/threshold",
        "Test-DbaSufficientDiskSpace: Validates disk space availability",
        "Initialize-Directories: Creates data/log directories if needed",
        "Enable-QueryStore: Configures Query Store (SQL 2016+)",
        "Write-Log: Centralized logging with timestamps and levels"
    ])
    
    slide = create_content_slide(prs, "⚙️ Configuration Parameters")
    
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    tf = textbox.text_frame
    tf.word_wrap = True
    
    config_text = """SqlInstance: Target SQL Server instance name
Database.Name: Database name to create
Database.DataDrive: Drive letter for data files (e.g., "G")
Database.LogDrive: Drive letter for log files (e.g., "L")
Database.ExpectedDatabaseSize: Expected size - determines file count
FileSizes.DataSize: Initial size per data file (e.g., "200MB")
FileSizes.DataGrowth: Auto-growth increment for data files
FileSizes.LogSize: Initial transaction log size
FileSizes.LogGrowth: Auto-growth increment for log file
FileSizes.FileSizeThreshold: Max size per file for calculation
LogFile: Path to log file for operation logging"""
    
    for line in config_text.split('\n'):
        if tf.text:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        
        p.text = line
        p.font.size = Pt(13)
        p.space_before = Pt(4)
        
        if ':' in line:
            param = line.split(':')[0]
            p.text = f"• {line}"
            p.font.name = "Courier New"
    
    slide = create_content_slide(prs, "🔄 Calculation Logic")
    
    box1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(2),
        Inches(8), Inches(1.2)
    )
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(155, 89, 182)
    box1.line.color.rgb = RGBColor(255, 255, 255)
    box1.line.width = Pt(2)
    
    tf = box1.text_frame
    tf.text = "File Count Calculation Formula"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf.vertical_anchor = 1
    
    box2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(3.5),
        Inches(7), Inches(1)
    )
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(236, 240, 241)
    box2.line.color.rgb = RGBColor(52, 73, 94)
    box2.line.width = Pt(2)
    
    tf = box2.text_frame
    tf.text = "if (ExpectedSize > Threshold)\n  numberOfFiles = Ceiling(ExpectedSize / Threshold)\nelse\n  numberOfFiles = 1"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.name = "Courier New"
    tf.paragraphs[0].font.bold = True
    tf.vertical_anchor = 1
    
    add_bullet_points(slide, Inches(1.5), Inches(5), Inches(7), Inches(1.5), [
        "Example 1: 5GB / 10GB → 5GB ≤ 10GB → 1 file",
        "Example 2: 50GB / 10GB → 50GB > 10GB → Ceiling(5) = 5 files",
        "Example 3: 100GB / 10GB → 100GB > 10GB → Ceiling(10) = 10 files (capped at 8)",
        "All files are created in the PRIMARY filegroup"
    ])
    
    slide = create_content_slide(prs, "💾 Disk Space Validation Logic")
    
    box1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(2),
        Inches(8), Inches(1.2)
    )
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(230, 126, 34)
    box1.line.color.rgb = RGBColor(255, 255, 255)
    box1.line.width = Pt(2)
    
    tf = box1.text_frame
    tf.text = "Space Requirement Calculation"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf.vertical_anchor = 1
    
    add_bullet_points(slide, Inches(1), Inches(3.5), Inches(8), Inches(3), [
        "1. Calculate file count from ExpectedDatabaseSize",
        "2. Calculate data space: FileCount × DataSize",
        "3. Add safety margin: RequiredSpace × 1.10 (10% buffer)",
        "4. Check data drive: AvailableSpace ≥ RequiredSpace + Margin",
        "5. Calculate log space: LogSize × 1.10",
        "6. Check log drive: AvailableSpace ≥ LogSize + Margin",
        "7. If same drive: Check combined requirement",
        "8. Fail early with clear error if insufficient space"
    ])
    
    slide = create_content_slide(prs, "🧪 Testing Coverage")
    
    test_areas = [
        ("Unit Tests", [
            "Size string conversion",
            "File count calculation",
            "Directory initialization",
            "Logging functionality"
        ], RGBColor(52, 152, 219)),
        ("Integration Tests", [
            "End-to-end workflow",
            "Configuration loading",
            "Error handling",
            "Database creation"
        ], RGBColor(46, 204, 113)),
        ("Edge Cases", [
            "Invalid inputs",
            "Insufficient space",
            "Same drive scenarios",
            "Maximum file limits"
        ], RGBColor(230, 126, 34))
    ]
    
    left_start = Inches(0.7)
    top_pos = Inches(2.2)
    box_width = Inches(2.8)
    box_height = Inches(3.5)
    spacing = Inches(0.3)
    
    for i, (category, items, color) in enumerate(test_areas):
        left_pos = left_start + (i * (box_width + spacing))
        
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left_pos, top_pos,
            box_width, box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = RGBColor(255, 255, 255)
        box.line.width = Pt(2)
        
        tf = box.text_frame
        tf.clear()
        
        p1 = tf.paragraphs[0]
        p1.text = category
        p1.alignment = PP_ALIGN.CENTER
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(255, 255, 255)
        
        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_before = Pt(6)
        
        tf.vertical_anchor = 1
    
    slide = create_content_slide(prs, "📋 Error Handling")
    add_bullet_points(slide, Inches(1), Inches(2), Inches(8), Inches(4.5), [
        "✅ Validates SQL Server connection before operations",
        "✅ Checks configuration file existence and format",
        "✅ Validates drive letters and paths",
        "✅ Verifies disk space before database creation",
        "✅ Handles existing database scenarios gracefully",
        "✅ Provides detailed error messages with context",
        "✅ Logs all operations for troubleshooting",
        "✅ Uses try-catch blocks for robust error handling"
    ])
    
    slide = create_content_slide(prs, "🔒 Best Practices Implemented")
    add_bullet_points(slide, Inches(1), Inches(2), Inches(8), Inches(4.5), [
        "📊 Caps data files at 8 (SQL Server recommendation)",
        "📁 Separates data and log files on different drives",
        "💾 Pre-validates disk space to prevent failures",
        "📝 Comprehensive logging for audit trail",
        "⚙️ Uses Query Store for performance monitoring (SQL 2016+)",
        "🔧 Sets database owner to 'sa' for consistency",
        "📚 Comment-based help for all functions",
        "🧪 Test coverage ensures reliability"
    ])
    
    prs.save('/home/ubuntu/sqlserver-databasescripts/SQL_Server_Technical_Documentation.pptx')
    print("✅ Created: SQL_Server_Technical_Documentation.pptx")

def create_user_guide_presentation():
    """Create user guide presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    create_title_slide(
        prs,
        "User Guide",
        "Step-by-Step Instructions for SQL Server Database Scripts"
    )
    
    slide = create_content_slide(prs, "📦 Prerequisites")
    add_bullet_points(slide, Inches(1), Inches(2), Inches(8), Inches(4.5), [
        "✅ PowerShell 5.1 or higher",
        "✅ SQL Server 2012 or higher (2016+ for Query Store)",
        "✅ dbatools PowerShell module",
        "✅ Appropriate SQL Server permissions (CREATE DATABASE)",
        "✅ Write access to data and log drive directories",
        "✅ Network connectivity to SQL Server instance"
    ])
    
    slide = create_content_slide(prs, "🚀 Step 1: Install dbatools")
    
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(2.5),
        Inches(7), Inches(3)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(44, 62, 80)
    box.line.color.rgb = RGBColor(52, 152, 219)
    box.line.width = Pt(3)
    
    tf = box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "PowerShell Command:"
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(52, 152, 219)
    
    p2 = tf.add_paragraph()
    p2.text = "\nInstall-Module -Name dbatools -Scope CurrentUser -Force -AllowClobber"
    p2.font.size = Pt(16)
    p2.font.name = "Courier New"
    p2.font.color.rgb = RGBColor(46, 204, 113)
    p2.space_before = Pt(12)
    
    p3 = tf.add_paragraph()
    p3.text = "\n\nThis installs the dbatools module required for SQL Server operations"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(149, 165, 166)
    p3.space_before = Pt(12)
    
    slide = create_content_slide(prs, "⚙️ Step 2: Configure DatabaseConfig.psd1")
    
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(2),
        Inches(8.4), Inches(4.5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(44, 62, 80)
    box.line.color.rgb = RGBColor(241, 196, 15)
    box.line.width = Pt(3)
    
    tf = box.text_frame
    tf.word_wrap = True
    
    config_example = """@{
    SqlInstance = "YourServerName"
    Database = @{
        Name = "MyDatabase"
        DataDrive = "G"
        LogDrive = "L"
        ExpectedDatabaseSize = "50GB"  # Results in 5 files (50GB / 10GB)
    }
    FileSizes = @{
        DataSize = "200MB"
        DataGrowth = "100MB"
        LogSize = "100MB"
        LogGrowth = "100MB"
        FileSizeThreshold = "10GB"
    }
    LogFile = "DatabaseCreation.log"
}"""
    
    p = tf.paragraphs[0]
    p.text = config_example
    p.font.size = Pt(13)
    p.font.name = "Courier New"
    p.font.color.rgb = RGBColor(46, 204, 113)
    
    slide = create_content_slide(prs, "▶️ Step 3: Run the Script")
    
    add_bullet_points(slide, Inches(1), Inches(2), Inches(8), Inches(1.5), [
        "Navigate to the script directory in PowerShell",
        "Execute the main script with configuration path"
    ])
    
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(4),
        Inches(7), Inches(2)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(44, 62, 80)
    box.line.color.rgb = RGBColor(46, 204, 113)
    box.line.width = Pt(3)
    
    tf = box.text_frame
    
    p1 = tf.paragraphs[0]
    p1.text = "cd C:\\path\\to\\sqlserver-databasescripts"
    p1.font.size = Pt(15)
    p1.font.name = "Courier New"
    p1.font.color.rgb = RGBColor(149, 165, 166)
    
    p2 = tf.add_paragraph()
    p2.text = "\n.\\SQLDatabaseCreation\\Invoke-DatabaseCreation.ps1 `\n  -ConfigPath .\\SQLDatabaseCreation\\DatabaseConfig.psd1"
    p2.font.size = Pt(15)
    p2.font.name = "Courier New"
    p2.font.color.rgb = RGBColor(46, 204, 113)
    p2.space_before = Pt(12)
    
    slide = create_content_slide(prs, "✅ What Happens Next")
    
    steps = [
        ("1", "Connection", "Validates SQL Server connection", RGBColor(52, 152, 219)),
        ("2", "Directories", "Creates data/log directories if needed", RGBColor(46, 204, 113)),
        ("3", "Calculation", "Calculates optimal file count", RGBColor(155, 89, 182)),
        ("4", "Validation", "Checks disk space availability", RGBColor(230, 126, 34)),
        ("5", "Creation", "Creates database with files", RGBColor(231, 76, 60)),
        ("6", "Configuration", "Sets owner and Query Store", RGBColor(26, 188, 156))
    ]
    
    left_start = Inches(0.5)
    top_start = Inches(2)
    box_width = Inches(3)
    box_height = Inches(0.7)
    
    for i, (num, title, desc, color) in enumerate(steps):
        row = i // 2
        col = i % 2
        
        left_pos = left_start + (col * (box_width + Inches(0.2)))
        top_pos = top_start + (row * (box_height + Inches(0.15)))
        
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left_pos, top_pos,
            box_width, box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = RGBColor(255, 255, 255)
        box.line.width = Pt(2)
        
        tf = box.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = f"{num}. {title}: {desc}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        tf.vertical_anchor = 1
    
    slide = create_content_slide(prs, "🔍 Monitoring Progress")
    add_bullet_points(slide, Inches(1), Inches(2), Inches(8), Inches(4.5), [
        "📊 Watch console output for real-time status updates",
        "📝 Check the log file for detailed operation information",
        "✅ Look for \"Success\" messages confirming each step",
        "⚠️ Any errors will be clearly displayed with context",
        "🔍 Use -Verbose flag for additional diagnostic information",
        "🧪 Use -WhatIf to preview without making changes"
    ])
    
    slide = create_content_slide(prs, "🛠️ Troubleshooting")
    
    issues = [
        ("Connection Failed", [
            "Verify SQL Server name",
            "Check network connectivity",
            "Confirm firewall rules"
        ]),
        ("Insufficient Space", [
            "Free up disk space",
            "Reduce file sizes",
            "Use different drives"
        ]),
        ("Permission Denied", [
            "Check SQL permissions",
            "Verify drive access",
            "Run as administrator"
        ])
    ]
    
    left_start = Inches(0.7)
    top_pos = Inches(2.2)
    box_width = Inches(2.8)
    box_height = Inches(3.5)
    spacing = Inches(0.3)
    
    colors = [RGBColor(231, 76, 60), RGBColor(230, 126, 34), RGBColor(155, 89, 182)]
    
    for i, (issue, solutions) in enumerate(issues):
        left_pos = left_start + (i * (box_width + spacing))
        
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left_pos, top_pos,
            box_width, box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i]
        box.line.color.rgb = RGBColor(255, 255, 255)
        box.line.width = Pt(2)
        
        tf = box.text_frame
        tf.clear()
        
        p1 = tf.paragraphs[0]
        p1.text = f"⚠️ {issue}"
        p1.alignment = PP_ALIGN.CENTER
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(255, 255, 255)
        
        p2 = tf.add_paragraph()
        p2.text = "\nSolutions:"
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.space_before = Pt(12)
        
        for solution in solutions:
            p = tf.add_paragraph()
            p.text = f"• {solution}"
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_before = Pt(4)
        
        tf.vertical_anchor = 1
    
    slide = create_content_slide(prs, "📚 Additional Resources")
    add_bullet_points(slide, Inches(1), Inches(2), Inches(8), Inches(4.5), [
        "📖 README.md: Complete documentation with examples",
        "🧪 Run Pester tests: Invoke-Pester -Path .\\Tests\\",
        "💡 Use Get-Help: Get-Help <FunctionName> -Full",
        "🔗 dbatools docs: https://dbatools.io/",
        "🔗 SQL Server best practices: Microsoft Learn",
        "💬 GitHub Issues: Report bugs or request features"
    ])
    
    slide = create_content_slide(prs, "✨ Tips & Best Practices")
    add_bullet_points(slide, Inches(1), Inches(2), Inches(8), Inches(4.5), [
        "💡 Use ExpectedDatabaseSize for automatic file calculation",
        "🔍 Always test with -WhatIf first in production",
        "📊 Keep data and log files on separate drives",
        "💾 Ensure 10-20% extra space beyond calculated requirement",
        "📝 Review logs after each operation",
        "🧪 Run Pester tests before deploying configuration changes",
        "🔄 Consider setting up as a scheduled task for automation",
        "📚 Document your specific configurations for team reference"
    ])
    
    prs.save('/home/ubuntu/sqlserver-databasescripts/SQL_Server_User_Guide.pptx')
    print("✅ Created: SQL_Server_User_Guide.pptx")

if __name__ == "__main__":
    print("🎨 Generating PowerPoint presentations...")
    print()
    
    create_overview_presentation()
    create_technical_presentation()
    create_user_guide_presentation()
    
    print()
    print("🎉 All presentations created successfully!")
    print()
    print("📁 Files created:")
    print("   1. SQL_Server_Database_Scripts_Overview.pptx")
    print("   2. SQL_Server_Technical_Documentation.pptx")
    print("   3. SQL_Server_User_Guide.pptx")
